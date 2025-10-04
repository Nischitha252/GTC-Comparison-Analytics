import os
import io
import requests
import subprocess
import xlrd
import pandas as pd
from pptx import Presentation
from io import BytesIO
from azure.storage.blob import BlobServiceClient
from tempfile import NamedTemporaryFile
from .Document_Processing import DocumentIntelligenceLoader
from azure.storage.blob import BlobServiceClient
import docx2txt
from .Document_Processing import BlobStorageProcessor
from flask import jsonify
from dotenv import load_dotenv
from urllib.parse import urlparse

load_dotenv()

AZURE_BLOB_STORAGE_CONNECTION_STRING = os.environ.get('AZURE_BLOB_STORAGE_CONNECTION_STRING')
blobstorageprocessor=BlobStorageProcessor()
documentintelligence=DocumentIntelligenceLoader()

def parse_blob_url(blob_url):
    parsed_url = urlparse(blob_url)
    path_parts = parsed_url.path.lstrip('/').split('/')
    container_name = path_parts[0]
    blob_name = '/'.join(path_parts[1:])
    return container_name, blob_name

def process_pdf_docx(blob_url):
    try:
        container_name, blob_name = parse_blob_url(blob_url)
        blob_url=blobstorageprocessor.get_blob_sas_url(container_name,blob_name)
        result = documentintelligence.analyze_document_word(blob_url)
        text=result.content
        return text
    except Exception as e:
        return jsonify({'error': f'Error reading PDF file: {str(e)}'}), 500

def process_doc(blob_url):
    container_name, blob_name = parse_blob_url(blob_url)
    blob_url=blobstorageprocessor.get_blob_sas_url(container_name,blob_name)
    response = requests.get(blob_url)
    doc_content = io.BytesIO(response.content)
    with NamedTemporaryFile(suffix='.doc', delete=False) as temp_doc:
        temp_doc.write(doc_content.getvalue())
        temp_doc_path = temp_doc.name

    output_dir = os.path.dirname(temp_doc_path)
    subprocess.call(['lowriter', '--headless', '--convert-to', 'docx', '--outdir', output_dir, temp_doc_path])
    temp_docx_path = temp_doc_path.replace('.doc', '.docx')
    text = docx2txt.process(temp_docx_path)

    os.remove(temp_doc_path)
    os.remove(temp_docx_path)
    return text

def process_xls(stream):
    workbook = xlrd.open_workbook(file_contents=stream.readall())
    all_sheets_content = ""
    for sheet_index in range(workbook.nsheets):
        sheet = workbook.sheet_by_index(sheet_index)
        sheet_content = f"Sheet: {sheet.name}\n"
        for row in range(sheet.nrows):
            row_values = []
            for col in range(sheet.ncols):
                cell_value = sheet.cell_value(row, col)
                if sheet.cell_type(row, col) == xlrd.XL_CELL_DATE:
                    cell_value = xlrd.xldate.xldate_as_datetime(cell_value, workbook.datemode).strftime('%Y-%m-%d %H:%M:%S')
                row_values.append(str(cell_value))
            sheet_content += "\t".join(row_values) + "\n"

        all_sheets_content += sheet_content + "\n\n"

    text = all_sheets_content
    return text

def process_xlsx(stream):
    excel_df = pd.read_excel(BytesIO(stream.readall()), sheet_name=None, engine='openpyxl')
    all_sheets_content = ""
    for sheet_name, df in excel_df.items():
        all_sheets_content += f"Sheet: {sheet_name}\n{df.to_csv(index=False)}"
    text = all_sheets_content
    return text

def process_csv(stream):
    csv_df = pd.read_csv(BytesIO(stream.readall()))
    text = csv_df.to_csv(index=False)
    return text

def process_pptx(stream):
    presentation = Presentation(BytesIO(stream.readall()))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def process_file(blob_url, file_format, stream=None):
    if file_format == '.pdf':
        return process_pdf_docx(blob_url)
    elif file_format == '.docx':
        return process_pdf_docx(blob_url)
    elif file_format == '.doc':
        return process_doc(blob_url)
    elif file_format == '.xls':
        return process_xls(stream)
    elif file_format == '.xlsx':
        return process_xlsx(stream)
    elif file_format == '.csv':
        return process_csv(stream)
    elif file_format == '.pptx':
        return process_pptx(stream)
    else:
        raise ValueError('Unsupported file format')
